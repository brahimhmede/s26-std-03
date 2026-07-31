"""PowerPoint renderer for normalized Futurdata disassembly guides.

The renderer contains only presentation concerns: slide layout, typography,
images, tables and pagination. It does not parse the original graph and it does
not mutate the JSON source. Non-obvious decisions are documented near the code,
especially branch-source resolution, image scaling, action pagination and the
separation between detailed and compact step layouts.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .base import Exporter
from ..models import Component, Step, WizardDocument
from ..options import ExportOptions
from ..utils import chunks, fit_rect, image_dimensions, normalize_text, resolve_image


class PPTXExporter(Exporter):
    format_id = "pptx"
    display_name = "PowerPoint 2007+"
    extension = ".pptx"

    WIDTH = 13.333
    HEIGHT = 7.5

    NAVY = "132238"
    NAVY_2 = "1E3450"
    TEAL = "19A7A0"
    SKY = "DFF5F4"
    WHITE = "FFFFFF"
    PAPER = "F7F9FC"
    INK = "14213D"
    MUTED = "64748B"
    LINE = "D9E2EC"
    WARNING = "F59E0B"
    ERROR = "DC2626"
    SUCCESS = "16A34A"

    def export(self, document: WizardDocument, output_path: Path, options: ExportOptions) -> Path:
        """Render ``document`` to an editable ``.pptx`` file.

        The method deliberately follows a linear pipeline: normalize options,
        select the requested action groups, generate front matter, generate
        operation slides, add consistent footers, and finally save atomically
        through python-pptx.
        """

        # Normalization protects the renderer from invalid GUI values such as
        # zero groups per slide or an inverted start/end range.
        options = options.normalized()

        # Step filtering happens before any slide is created, so all counters
        # and the title slide accurately describe the exported subset.
        selected_steps = document.selected_steps(
            options.start_step,
            options.end_step,
            options.max_action_groups,
        )
        if not selected_steps:
            raise ValueError("The selected export range contains no disassembly steps.")

        # Start from a blank presentation and force a 16:9 canvas. The deck
        # remains a normal editable PowerPoint rather than a flattened image.
        prs = Presentation()
        prs.slide_width = Inches(self.WIDTH)
        prs.slide_height = Inches(self.HEIGHT)

        # Front matter is controlled independently so the same renderer can
        # produce a complete report or a compact sequence of operation slides.
        if options.include_title:
            self._title_slide(prs, document, selected_steps, options)
        if options.include_overview:
            self._overview_slide(prs, document, selected_steps)
        if options.include_warnings:
            self._warning_slides(prs, document)
        if options.include_tools_summary or options.include_safety_summary:
            self._tools_safety_slide(prs, document, options)
        if options.include_bom:
            self._bom_slides(prs, document, options)

        # One group per slide uses the detailed layout. Two to four groups
        # use compact cards, implementing the URS choice of action groups/page.
        for group in chunks(selected_steps, options.groups_per_slide):
            if options.groups_per_slide == 1:
                self._detailed_step_slides(prs, document, group[0], options)
            else:
                self._grouped_step_slide(prs, document, group, options)

        if options.include_closing:
            self._closing_slide(prs, document, selected_steps)

        # Footers are added only after every slide exists; this guarantees a
        # correct total-page value even when warnings and BoM require pagination.
        for number, slide in enumerate(prs.slides, start=1):
            self._footer(slide, document.product.name, number, len(prs.slides))

        destination = self.ensure_extension(Path(output_path).expanduser().resolve())
        destination.parent.mkdir(parents=True, exist_ok=True)
        try:
            prs.save(str(destination))
        except PermissionError as exc:
            raise PermissionError(
                f"Cannot save '{destination}'. Close the file if it is open in PowerPoint and try again."
            ) from exc
        return destination

    # ---------- General drawing helpers ----------

    @staticmethod
    def _rgb(value: str) -> RGBColor:
        return RGBColor.from_string(value)

    def _blank(self, prs: Presentation):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide, self.PAPER)
        return slide

    def _set_background(self, slide, color: str) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = self._rgb(color)

    def _rect(self, slide, x, y, w, h, fill, line=None, radius=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill)
        shape.line.color.rgb = self._rgb(line or fill)
        if radius:
            shape.adjustments[0] = 0.08
        return shape

    def _text(
        self,
        slide,
        value,
        x,
        y,
        w,
        h,
        size=16,
        bold=False,
        color=None,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        font="Aptos",
        margin=0.03,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = str(value)
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color or self.INK)
        return box

    def _section_title(self, slide, eyebrow: str, title: str, subtitle: str | None = None):
        self._text(slide, eyebrow.upper(), 0.65, 0.36, 5.8, 0.28, 10, True, self.TEAL)
        self._text(slide, title, 0.65, 0.7, 12.0, 0.55, 27, True, self.INK)
        if subtitle:
            self._text(slide, subtitle, 0.65, 1.3, 12.0, 0.45, 12, False, self.MUTED)

    def _footer(self, slide, product_name: str, number: int, total: int):
        self._text(slide, product_name, 0.65, 7.16, 9.5, 0.18, 8, False, self.MUTED)
        self._text(slide, f"{number} / {total}", 11.85, 7.16, 0.8, 0.18, 8, True, self.MUTED, PP_ALIGN.RIGHT)

    def _picture(self, slide, value: str | None, source_dir: str | None, x, y, w, h) -> bool:
        """Insert an image inside a bounding box while preserving aspect ratio.

        Invalid or missing images are treated as optional metadata: the method
        returns ``False`` and the caller can draw a placeholder instead of
        aborting the entire export.
        """

        source = resolve_image(value, source_dir)
        if source is None:
            return False
        try:
            width, height = image_dimensions(source)
            fitted_w, fitted_h = fit_rect(width, height, w, h)
            left = x + (w - fitted_w) / 2
            top = y + (h - fitted_h) / 2
            if isinstance(source, io.BytesIO):
                source.seek(0)
            slide.shapes.add_picture(source, Inches(left), Inches(top), Inches(fitted_w), Inches(fitted_h))
            return True
        except Exception:
            return False

    # ---------- Front matter ----------

    def _title_slide(self, prs, document, selected_steps, options):
        slide = self._blank(prs)
        self._set_background(slide, self.NAVY)
        self._rect(slide, 0, 0, 0.18, 7.48, self.TEAL, radius=False)
        self._text(slide, "FUTURDATA THESIS", 0.85, 0.75, 5.0, 0.35, 12, True, self.TEAL)
        self._text(slide, "Disassembly Guide", 0.85, 1.35, 7.5, 0.8, 38, True, self.WHITE)
        self._text(slide, document.product.name, 0.85, 2.25, 7.5, 1.15, 26, False, self.WHITE)
        summary = (
            f"PPTX export • {len(selected_steps)} action group(s) • "
            f"depth: {document.depth_mode} • schema {document.schema_version}"
        )
        self._text(slide, summary, 0.85, 3.7, 7.5, 0.6, 14, False, "C7D4E3")
        self._rect(slide, 0.85, 5.35, 3.8, 0.8, self.NAVY_2, self.NAVY_2)
        self._text(slide, f"Nominal weight  {document.product.weight_label}", 1.1, 5.58, 3.25, 0.3, 13, True, self.WHITE)
        image_present = self._picture(slide, document.product.image, document.source_dir, 8.9, 0.9, 3.7, 5.8)
        if not image_present:
            self._rect(slide, 9.15, 1.25, 3.15, 4.85, self.NAVY_2, "34516F")
            self._text(slide, "PRODUCT\nIMAGE", 9.55, 2.95, 2.35, 1.0, 20, True, "7B92AB", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    def _overview_slide(self, prs, document, selected_steps):
        slide = self._blank(prs)
        self._section_title(slide, "Project overview", "Export scope and model summary")
        cards = [
            ("PRODUCT ID", str(document.product.node_id if document.product.node_id is not None else "—")),
            ("ACTION GROUPS", str(len(selected_steps))),
            ("BOM COMPONENTS", str(len(document.bill_of_materials))),
            ("VALIDATION ISSUES", str(len(document.validation_issues or document.warnings))),
        ]
        for index, (label, value) in enumerate(cards):
            x = 0.65 + index * 3.08
            self._rect(slide, x, 1.95, 2.75, 1.25, self.WHITE, self.LINE)
            self._text(slide, label, x + 0.22, 2.16, 2.3, 0.25, 9, True, self.TEAL)
            self._text(slide, value, x + 0.22, 2.5, 2.3, 0.45, 24, True, self.INK)

        self._rect(slide, 0.65, 3.55, 5.95, 2.7, self.WHITE, self.LINE)
        self._text(slide, "PRODUCT DATA", 0.95, 3.83, 2.5, 0.25, 10, True, self.TEAL)
        product_lines = [
            f"Name: {document.product.name}",
            f"Weight: {document.product.weight_label}",
            f"Material: {document.lookup_label('materials', document.product.material)}",
            f"Color: {document.lookup_label('colors', document.product.color)}",
            f"Depth mode: {document.depth_mode}",
        ]
        if document.keep_whole_ids:
            product_lines.append(f"Keep-whole IDs: {', '.join(map(str, document.keep_whole_ids))}")
        self._text(slide, "\n".join(product_lines), 0.95, 4.18, 5.3, 1.75, 14, False, self.INK)

        self._rect(slide, 6.85, 3.55, 5.82, 2.7, self.WHITE, self.LINE)
        self._text(slide, "EXPORT BEHAVIOUR", 7.15, 3.83, 2.8, 0.25, 10, True, self.TEAL)
        behavior = [
            "• Parent operation and atomic actions stay grouped.",
            "• Output components and remaining assembly are shown.",
            "• Text and images remain editable in PowerPoint.",
            "• Unknown branch sources are never guessed.",
            "• Source JSON is opened read-only and never modified.",
        ]
        self._text(slide, "\n".join(behavior), 7.15, 4.18, 5.05, 1.75, 13, False, self.INK)

    def _warning_slides(self, prs, document):
        """Paginate validation issues, six per slide, with severity accents."""

        issues = document.validation_issues or [
            type("Issue", (), {
                "severity": warning.severity,
                "rule": warning.rule,
                "message": warning.message,
                "location": warning.location,
            })
            for warning in document.warnings
        ]
        if not issues:
            return
        for page, group in enumerate(chunks(list(issues), 6), start=1):
            slide = self._blank(prs)
            suffix = f" — {page}" if len(issues) > 6 else ""
            self._section_title(slide, "Validation", f"Model warnings and errors{suffix}", "Each issue is localized to the corresponding JSON element when possible.")
            y = 1.95
            for issue in group:
                severity = str(issue.severity).lower()
                accent = self.ERROR if severity == "error" else self.WARNING
                self._rect(slide, 0.65, y, 12.0, 0.72, self.WHITE, self.LINE)
                self._rect(slide, 0.65, y, 0.12, 0.72, accent, accent, radius=False)
                self._text(slide, f"{severity.upper()} · {issue.rule}", 0.95, y + 0.12, 3.1, 0.22, 10, True, accent)
                location = f"Location: {issue.location}" if getattr(issue, "location", "") else ""
                self._text(slide, normalize_text(issue.message), 3.25, y + 0.1, 7.9, 0.32, 11, False, self.INK)
                self._text(slide, location, 3.25, y + 0.43, 8.7, 0.18, 8, False, self.MUTED)
                y += 0.82

    def _tools_safety_slide(self, prs, document, options):
        tools = document.tools_summary if options.include_tools_summary else []
        safety = document.safety_summary if options.include_safety_summary else []
        slide = self._blank(prs)
        self._section_title(slide, "Preparation", "Tools and safety notices", "Aggregated from step-level and atomic-action metadata.")
        self._rect(slide, 0.65, 1.9, 5.85, 4.85, self.WHITE, self.LINE)
        self._text(slide, "TOOLS REQUIRED", 0.95, 2.2, 3.3, 0.25, 11, True, self.TEAL)
        tool_lines = [f"• {tool}" for tool in tools] or ["• No tools are specified in this model."]
        self._text(slide, "\n".join(tool_lines[:15]), 0.95, 2.65, 5.05, 3.65, 14, False, self.INK)

        self._rect(slide, 6.8, 1.9, 5.85, 4.85, self.WHITE, self.LINE)
        self._text(slide, "SAFETY NOTICES", 7.1, 2.2, 3.3, 0.25, 11, True, self.WARNING)
        safety_lines = [f"• {notice}" for notice in safety] or ["• No safety notices are specified in this model."]
        self._text(slide, "\n".join(safety_lines[:15]), 7.1, 2.65, 5.05, 3.65, 14, False, self.INK)

    def _bom_slides(self, prs, document, options):
        """Create editable BoM tables and split long lists across slides."""

        bom = document.bill_of_materials
        if not bom:
            return
        groups = list(chunks(bom, options.bom_rows_per_slide))
        for page, group in enumerate(groups, start=1):
            slide = self._blank(prs)
            self._section_title(slide, "Bill of Materials", f"Recovered components — page {page} of {len(groups)}")
            x, y, w, h = 0.55, 1.72, 12.25, 5.15
            rows = len(group) + 1
            cols = 7
            table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
            widths = [0.7, 4.1, 1.15, 1.35, 1.05, 1.65, 1.8]
            for index, width in enumerate(widths):
                table.columns[index].width = Inches(width)
            headers = ["ID", "Component", "Weight", "Material", "Color", "Quality", "Destination"]
            for col, header in enumerate(headers):
                cell = table.cell(0, col)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb(self.NAVY)
                self._format_cell(cell, 10, True, self.WHITE)
            for row, component in enumerate(group, start=1):
                values = [
                    component.node_id if component.node_id is not None else "—",
                    component.name,
                    component.weight_label,
                    document.lookup_label("materials", component.material),
                    document.lookup_label("colors", component.color),
                    component.quality or "—",
                    component.destination or "—",
                ]
                for col, value in enumerate(values):
                    cell = table.cell(row, col)
                    cell.text = str(value)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._rgb(self.WHITE if row % 2 else "F0F5FA")
                    self._format_cell(cell, 9 if col != 1 else 10, col == 1, self.INK)

    def _format_cell(self, cell, size, bold, color):
        cell.margin_left = Inches(0.07)
        cell.margin_right = Inches(0.07)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        frame = cell.text_frame
        frame.word_wrap = True
        for paragraph in frame.paragraphs:
            paragraph.space_after = Pt(0)
            for run in paragraph.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = self._rgb(color)

    # ---------- Disassembly action slides ----------

    def _detailed_step_slides(self, prs, document, step, options):
        """Render one parent operation, paginating long atomic-action lists."""

        # Six action rows fit safely in the available content region. Longer
        # lists are continued on additional slides instead of shrinking text.
        action_lines = [normalize_text(action.text) for action in step.actions if normalize_text(action.text)]
        if not action_lines:
            action_lines = [normalize_text(step.operation)]
        action_pages = list(chunks(action_lines, 6))
        for page, page_actions in enumerate(action_pages, start=1):
            slide = self._blank(prs)
            page_label = f" · part {page}/{len(action_pages)}" if len(action_pages) > 1 else ""
            self._section_title(
                slide,
                f"Action group {step.index}",
                f"{step.operation}{page_label}",
                "Atomic actions are grouped under their parent disassembly operation.",
            )
            # Source resolution is graph-safe. An explicit source is used when
            # available; otherwise only a genuine linear continuation is inferred.
            # A new branch is labelled unknown rather than incorrectly reset to
            # the product root.
            position = document.position_of(step)
            source, method = document.source_for_step(position)
            if source:
                source_text = f"Source: {source.name} ({source.weight_label})"
            else:
                source_text = "Source: not specified in JSON — new branch"
            method_text = {
                "explicit": "explicit source",
                "product_root": "product root",
                "previous_continuation": "previous continuation",
                "unspecified_branch": "unresolved branch",
            }[method]
            self._rect(slide, 0.65, 1.7, 12.0, 0.48, self.SKY, self.SKY)
            self._text(slide, source_text, 0.9, 1.82, 9.0, 0.22, 11, True, self.INK)
            self._text(slide, method_text, 10.35, 1.82, 1.95, 0.22, 9, True, self.TEAL, PP_ALIGN.RIGHT)

            images = self._step_images(step)
            image_count = len(images) if options.include_images else 0
            content_width = 7.45
            self._rect(slide, 0.65, 2.4, content_width, 4.35, self.WHITE, self.LINE)
            self._text(slide, "ATOMIC ACTIONS", 0.95, 2.68, 3.0, 0.25, 10, True, self.TEAL)
            y = 3.08
            for number, action in enumerate(page_actions, start=1 + (page - 1) * 6):
                self._rect(slide, 0.95, y, 0.38, 0.38, self.TEAL, self.TEAL)
                self._text(slide, number, 0.95, y + 0.02, 0.38, 0.28, 10, True, self.WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
                self._text(slide, action, 1.48, y - 0.02, content_width - 1.2, 0.52, 14, False, self.INK)
                y += 0.57

            if image_count:
                self._image_gallery(slide, images[:3], document.source_dir, 8.35, 2.4, 4.3, 2.05)
                output_y = 4.65
                output_h = 2.1
            else:
                output_y = 2.4
                output_h = 4.35
            right_x = 8.35
            right_w = 4.3
            self._rect(slide, right_x, output_y, right_w, output_h, self.WHITE, self.LINE)
            self._step_result_panel(slide, document, step, right_x, output_y, right_w, output_h, options)

    def _step_images(self, step: Step) -> list[tuple[str, str]]:
        """Collect and de-duplicate images from operation, actions and outputs."""

        values: list[tuple[str, str]] = []
        if step.image:
            values.append((step.image, "Operation"))
        for index, action in enumerate(step.actions, start=1):
            if action.image:
                values.append((action.image, f"Action {index}"))
        for component in step.outputs:
            if component.image:
                values.append((component.image, component.name))
        if step.continues_as and step.continues_as.image:
            values.append((step.continues_as.image, step.continues_as.name))
        deduped: list[tuple[str, str]] = []
        seen: set[str] = set()
        for value, caption in values:
            if value not in seen:
                seen.add(value)
                deduped.append((value, caption))
        return deduped

    def _image_gallery(self, slide, images, source_dir, x, y, w, h):
        self._rect(slide, x, y, w, h, self.WHITE, self.LINE)
        count = len(images)
        box_w = (w - 0.3 - 0.12 * (count - 1)) / count
        for index, (value, caption) in enumerate(images):
            left = x + 0.15 + index * (box_w + 0.12)
            if not self._picture(slide, value, source_dir, left, y + 0.15, box_w, h - 0.55):
                self._rect(slide, left, y + 0.15, box_w, h - 0.55, "EEF2F7", self.LINE)
            self._text(slide, caption, left, y + h - 0.34, box_w, 0.18, 8, True, self.MUTED, PP_ALIGN.CENTER)

    def _step_result_panel(self, slide, document, step, x, y, w, h, options):
        self._text(slide, "RESULT", x + 0.25, y + 0.22, 1.5, 0.22, 10, True, self.TEAL)
        outputs = step.outputs or []
        lines = [f"• {component.name} ({component.weight_label})" for component in outputs]
        if not lines:
            lines = ["• No separate output component encoded"]
        max_output_lines = 3 if h < 3 else 5
        self._text(slide, "\n".join(lines[:max_output_lines]), x + 0.25, y + 0.55, w - 0.5, 0.75 if h < 3 else 1.3, 10 if h < 3 else 11, False, self.INK)
        continuation = (
            f"{step.continues_as.name} ({step.continues_as.weight_label})"
            if step.continues_as else "End of this branch"
        )
        continuation_y = y + (1.35 if h < 3 else 2.0)
        self._text(slide, "CONTINUES AS", x + 0.25, continuation_y, w - 0.5, 0.2, 9, True, self.TEAL)
        self._text(slide, continuation, x + 0.25, continuation_y + 0.24, w - 0.5, 0.42, 10, False, self.INK)
        if h >= 3:
            tools = ", ".join(step.all_tools) if step.all_tools else "No tools specified"
            self._text(slide, "TOOLS", x + 0.25, y + 2.85, w - 0.5, 0.2, 9, True, self.TEAL)
            self._text(slide, tools, x + 0.25, y + 3.08, w - 0.5, 0.45, 10, False, self.INK)
        if options.include_component_details and outputs:
            component = outputs[0]
            details = []
            if component.material not in (None, ""):
                details.append(f"Material: {document.lookup_label('materials', component.material)}")
            if component.color not in (None, ""):
                details.append(f"Color: {document.lookup_label('colors', component.color)}")
            if component.quality:
                details.append(f"Quality: {component.quality}")
            if component.destination:
                details.append(f"Destination: {component.destination}")
            if details:
                self._text(slide, " · ".join(details), x + 0.25, y + h - 0.36, w - 0.5, 0.18, 8, False, self.MUTED)

    def _grouped_step_slide(self, prs, document, steps, options):
        """Render two to four parent operations as compact sequence cards."""

        slide = self._blank(prs)
        first, last = steps[0].index, steps[-1].index
        self._section_title(slide, "Disassembly sequence", f"Action groups {first}–{last}", f"{len(steps)} parent operations on this slide.")
        available_height = 5.15
        gap = 0.14
        card_h = (available_height - gap * (len(steps) - 1)) / len(steps)
        y = 1.75
        for step in steps:
            self._compact_step_card(slide, document, step, 0.65, y, 12.0, card_h)
            y += card_h + gap

    def _compact_step_card(self, slide, document, step, x, y, w, h):
        self._rect(slide, x, y, w, h, self.WHITE, self.LINE)
        self._rect(slide, x, y, 0.72, h, self.NAVY, self.NAVY, radius=False)
        self._text(slide, step.index, x, y + 0.12, 0.72, 0.35, 16, True, self.WHITE, PP_ALIGN.CENTER)
        self._text(slide, step.operation, x + 0.92, y + 0.12, 2.45, 0.38, 16, True, self.INK)
        position = document.position_of(step)
        source, _ = document.source_for_step(position)
        source_label = source.name if source else "Source not specified"
        self._text(slide, f"From: {source_label}", x + 0.92, y + 0.55, 2.55, 0.28, 9, False, self.MUTED)
        actions = "; ".join(normalize_text(action.text) for action in step.actions if normalize_text(action.text)) or normalize_text(step.operation)
        self._text(slide, actions, x + 3.5, y + 0.12, 4.05, h - 0.24, 10, False, self.INK)
        outputs = ", ".join(component.name for component in step.outputs) or "No separate output"
        self._text(slide, outputs, x + 7.72, y + 0.12, 2.3, h - 0.24, 10, False, self.INK)
        continuation = step.continues_as.name if step.continues_as else "End branch"
        self._text(slide, continuation, x + 10.18, y + 0.12, 1.55, h - 0.24, 9, False, self.MUTED)

    def _closing_slide(self, prs, document, selected_steps):
        """Add a final summary confirming that the deck remains editable."""

        slide = self._blank(prs)
        self._set_background(slide, self.NAVY)
        self._text(slide, "EXPORT COMPLETE", 0.85, 1.35, 5.2, 0.32, 12, True, self.TEAL)
        self._text(slide, f"{len(selected_steps)} action groups prepared", 0.85, 2.0, 10.5, 0.85, 36, True, self.WHITE)
        self._text(
            slide,
            "The PowerPoint remains fully editable: text, tables, images and component information can be revised after export.",
            0.85,
            3.25,
            10.6,
            1.0,
            18,
            False,
            self.WHITE,
        )
        self._rect(slide, 0.85, 5.25, 5.4, 0.85, self.NAVY_2, self.NAVY_2)
        self._text(slide, document.product.name, 1.1, 5.5, 4.9, 0.32, 13, True, self.TEAL)
