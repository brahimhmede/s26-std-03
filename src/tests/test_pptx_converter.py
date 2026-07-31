"""Automated regression tests for the Futurdata PPTX converter."""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest
from pptx import Presentation

# The requested project directory contains a hyphen. Add the directory itself
# to Python's import path, exactly as the shared GUI does.
REPOSITORY_ROOT = Path(__file__).resolve().parents[2]
CONVERTER_DIR = REPOSITORY_ROOT / "src" / "main" / "pptx-converter"
if str(CONVERTER_DIR) not in sys.path:
    sys.path.insert(0, str(CONVERTER_DIR))

from pptx_exporter import export_to_pptx  # noqa: E402
from pptx_core import PPTXExportEngine  # noqa: E402

SAMPLE = CONVERTER_DIR / "data" / "example_air_fryer.json"


def _presentation_text(path: Path) -> str:
    """Return all editable slide text and table-cell values."""

    values: list[str] = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                values.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def test_complete_export_contains_main_sections(tmp_path: Path) -> None:
    """The default export must contain the main report sections."""

    output = tmp_path / "complete.pptx"
    result = Path(export_to_pptx(SAMPLE, output))

    assert result.exists()
    assert result.suffix == ".pptx"

    presentation = Presentation(result)
    assert len(presentation.slides) >= 10

    text = _presentation_text(result).lower()
    assert "disassembly guide" in text
    assert "bill of materials" in text
    assert "remove the cover" in text
    assert "plastic cover (8 hooks)" in text


def test_bom_switch_removes_bom_slides(tmp_path: Path) -> None:
    """The GUI BoM switch must control both data and slide generation."""

    result = Path(
        export_to_pptx(
            SAMPLE,
            tmp_path / "without_bom.pptx",
            include_bom=False,
        )
    )
    assert result.exists()
    assert "bill of materials" not in _presentation_text(result).lower()


def test_step_range_and_grouping(tmp_path: Path) -> None:
    """Step filtering and action-groups-per-slide must be deterministic."""

    result = Path(
        export_to_pptx(
            SAMPLE,
            tmp_path / "partial.pptx",
            start_step=3,
            end_step=8,
            groups_per_slide=2,
            include_bom=False,
            include_tools_summary=False,
            include_safety_summary=False,
        )
    )

    text = _presentation_text(result)
    assert "Action groups 3–4" in text
    assert "Action groups 7–8" in text
    assert "Action groups 9–10" not in text


def test_existing_guide_can_be_reused(tmp_path: Path) -> None:
    """The GUI may pass its already-built guide to avoid duplicate parsing."""

    engine = PPTXExportEngine()
    guide = engine.load(SAMPLE, prefer_external_loader=False)

    result = Path(
        export_to_pptx(
            json_path=SAMPLE,
            guide=guide,
            output_path=tmp_path / "from_guide.pptx",
        )
    )

    assert result.exists()
    assert "Air fryer Philips HD9252" in _presentation_text(result)


def test_unknown_branch_source_is_not_invented(tmp_path: Path) -> None:
    """A new branch without explicit source must remain visibly unknown."""

    result = Path(
        export_to_pptx(
            SAMPLE,
            tmp_path / "branch.pptx",
            start_step=9,
            end_step=9,
            include_bom=False,
        )
    )
    assert "Source: not specified in JSON" in _presentation_text(result)


def test_image_dictionary_from_loader_is_supported(tmp_path: Path) -> None:
    """Loader image objects converted to dictionaries must not crash export."""

    payload = json.loads(SAMPLE.read_text(encoding="utf-8"))
    payload["product"]["image"] = {
        "path": "missing-product-image.png",
        "is_url": False,
    }
    first_action = next(
        action
        for step in payload["steps"]
        for action in step.get("actions", [])
    )
    first_action["image"] = {
        "path": "missing-action-image.png",
        "is_url": False,
    }
    source = tmp_path / "image_dict.json"
    source.write_text(json.dumps(payload), encoding="utf-8")

    result = Path(export_to_pptx(source, tmp_path / "image_dict.pptx"))
    assert result.exists()
    assert len(Presentation(result).slides) >= 1


def test_invalid_json_has_readable_error(tmp_path: Path) -> None:
    """Malformed JSON should report its line/column instead of a raw traceback."""

    source = tmp_path / "broken.json"
    source.write_text('{"product": ', encoding="utf-8")

    with pytest.raises(ValueError, match="not valid JSON"):
        export_to_pptx(source, tmp_path / "broken.pptx")


def test_output_extension_is_added(tmp_path: Path) -> None:
    """Callers may omit the extension without creating an extensionless file."""

    result = Path(export_to_pptx(SAMPLE, tmp_path / "guide_without_extension"))
    assert result.name == "guide_without_extension.pptx"
    assert result.exists()
